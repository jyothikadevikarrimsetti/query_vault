#!/usr/bin/env python3.12
"""Generate QueryVault Sales Deck (PowerPoint) for C-Suite / Hospital CIO audience."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x2A, 0x4A)
DARK_NAVY  = RGBColor(0x0F, 0x17, 0x2A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0x93, 0xC5, 0xFD)
BLUE       = RGBColor(0x3B, 0x82, 0xF6)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
RED        = RGBColor(0xEF, 0x44, 0x44)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
GRAY       = RGBColor(0x94, 0xA3, 0xB8)
DARK_GRAY  = RGBColor(0x64, 0x74, 0x8B)
LIGHT_NAVY = RGBColor(0x23, 0x3A, 0x60)

SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)

# ── Helpers ─────────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor = NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(
    slide, left, top, width, height, text,
    font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
    font_name="Calibri", line_spacing=1.2,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_box(
    slide, left, top, width, height, lines,
    font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
    font_name="Calibri", bullet=False, line_spacing=1.5,
):
    """lines: list of (text, optional_color, optional_bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, clr, bld = item, color, bold
        elif len(item) == 2:
            txt, clr = item
            bld = bold
        else:
            txt, clr, bld = item

        if bullet:
            txt = f"\u2022  {txt}"

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)

    return txBox


def add_rounded_rect(
    slide, left, top, width, height,
    fill_color=LIGHT_NAVY, text="", font_size=14, font_color=WHITE,
    bold=False, alignment=PP_ALIGN.CENTER,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = alignment
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = "Calibri"
    return shape


def add_arrow(slide, left, top, width=Inches(0.5), height=Inches(0.05), color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_divider(slide, top, color=BLUE, left=Inches(0.8), width=Inches(11.7)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_metric_card(slide, left, top, number, label, accent=BLUE):
    card_w, card_h = Inches(2.4), Inches(1.8)
    add_rounded_rect(slide, left, top, card_w, card_h, LIGHT_NAVY)
    add_text_box(slide, left, top + Inches(0.2), card_w, Inches(0.8),
                 str(number), font_size=40, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(1.05), card_w, Inches(0.5),
                 label, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)


# ── Slide builders ──────────────────────────────────────────

def slide_01_title(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_NAVY)

    # Decorative accent line
    add_divider(slide, Inches(2.3), BLUE, Inches(4.5), Inches(4.3))

    add_text_box(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(0.8),
                 "QUERYVAULT", font_size=48, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.0),
                 "AI-Powered Database Security\nfor Healthcare", font_size=32, color=WHITE,
                 bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.6),
                 "Zero-Trust Protection for Every AI-Generated Query",
                 font_size=18, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.5),
                 "Apollo Hospitals  |  AI Query Security Platform",
                 font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Shield icon placeholder
    add_rounded_rect(slide, Inches(6.1), Inches(5.8), Inches(1.1), Inches(1.1),
                     BLUE, "\U0001F6E1", font_size=36, font_color=WHITE)


def slide_02_problem(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "THE PROBLEM", font_size=14, color=BLUE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "AI Meets Sensitive Data: A New Attack Surface",
                 font_size=32, color=WHITE, bold=True)

    add_divider(slide, Inches(1.6))

    # Three problem cards
    cards = [
        ("\U0001F4A5", "LLM Hallucinations", "AI models fabricate table & column\nnames, generating invalid or\nunauthorized SQL queries", RED),
        ("\U0001F6A8", "Prompt Injection", "Attackers embed malicious instructions\nto manipulate AI into accessing\nrestricted data or dropping tables", AMBER),
        ("\U0001F4CB", "Regulatory Pressure", "HIPAA, GDPR, EU AI Act, SOX\ndemand auditable, controlled AI\naccess to patient & financial data", BLUE),
    ]

    for i, (icon, title, desc, accent) in enumerate(cards):
        x = Inches(0.8 + i * 4.1)
        add_rounded_rect(slide, x, Inches(2.0), Inches(3.6), Inches(3.2), LIGHT_NAVY)
        add_text_box(slide, x + Inches(0.3), Inches(2.2), Inches(3.0), Inches(0.6),
                     icon, font_size=36, color=accent, alignment=PP_ALIGN.LEFT)
        add_text_box(slide, x + Inches(0.3), Inches(2.9), Inches(3.0), Inches(0.5),
                     title, font_size=20, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.3), Inches(3.5), Inches(3.0), Inches(1.4),
                     desc, font_size=14, color=GRAY, line_spacing=1.4)

    # Stats bar
    add_rounded_rect(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.3), DARK_NAVY)

    stats = [
        ("$10.93M", "Avg. healthcare\nbreach cost (2023)"),
        ("725", "Healthcare breaches\nreported in 2023"),
        ("88%", "Of breaches involve\nhuman error + AI"),
        ("$1.76M", "Average savings with\nAI security controls"),
    ]
    for i, (num, lbl) in enumerate(stats):
        x = Inches(1.2 + i * 3.0)
        add_text_box(slide, x, Inches(5.7), Inches(2.0), Inches(0.6),
                     num, font_size=28, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(6.25), Inches(2.0), Inches(0.6),
                     lbl, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)


def slide_03_solution(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "THE SOLUTION", font_size=14, color=GREEN, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "QueryVault: Zero-Trust AI Query Security",
                 font_size=32, color=WHITE, bold=True)

    add_divider(slide, Inches(1.6), GREEN)

    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(10), Inches(0.8),
                 "A security orchestration layer that wraps around any NL-to-SQL AI pipeline,\n"
                 "enforcing identity, access control, and compliance on every single query.",
                 font_size=18, color=LIGHT_BLUE, line_spacing=1.5)

    # Architecture flow
    flow_items = [
        ("User Question", GRAY),
        ("Identity &\nThreat Scan", RED),
        ("Context\nMinimization", AMBER),
        ("AI Model\n(XenSQL)", BLUE),
        ("SQL\nValidation", AMBER),
        ("Secure\nExecution", GREEN),
        ("Audit &\nMonitoring", LIGHT_BLUE),
    ]

    y_flow = Inches(3.3)
    for i, (label, color) in enumerate(flow_items):
        x = Inches(0.5 + i * 1.8)
        add_rounded_rect(slide, x, y_flow, Inches(1.5), Inches(1.2),
                         LIGHT_NAVY, label, font_size=12, font_color=color, bold=True)
        if i < len(flow_items) - 1:
            add_arrow(slide, x + Inches(1.55), y_flow + Inches(0.55), Inches(0.25), Inches(0.1), DARK_GRAY)

    # Key benefits
    benefits = [
        ("Every query inspected", "5 sequential security zones, zero exceptions"),
        ("Least-privilege AI", "AI only sees schema the user's role permits"),
        ("Tamper-proof audit", "SHA-256 hash-chain links every event"),
        ("7 compliance frameworks", "HIPAA, GDPR, EU AI Act, SOX and more"),
    ]

    for i, (title, desc) in enumerate(benefits):
        x = Inches(0.8 + i * 3.1)
        y = Inches(5.2)
        add_text_box(slide, x, y, Inches(2.8), Inches(0.4),
                     f"\u2713  {title}", font_size=14, color=GREEN, bold=True)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.4), Inches(2.5), Inches(0.5),
                     desc, font_size=12, color=GRAY)


def slide_04_pipeline(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "ARCHITECTURE", font_size=14, color=BLUE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "5-Zone Security Pipeline",
                 font_size=32, color=WHITE, bold=True)

    add_divider(slide, Inches(1.6))

    zones = [
        ("ZONE 1", "PRE-MODEL", "Identity verification, injection\nscanning, behavioral analysis,\nthreat classification", RED, "\U0001F6E1"),
        ("ZONE 2", "MODEL BOUNDARY", "Context minimization, schema\nfiltering, least-privilege\ndata exposure to AI", AMBER, "\U0001F512"),
        ("ZONE 3", "POST-MODEL", "3-gate SQL validation,\nhallucination detection,\nauto-masking & rewriting", BLUE, "\u2714"),
        ("ZONE 4", "EXECUTION", "Circuit breaker, read-only\nenforcement, connection\npooling, query timeout", GREEN, "\u26A1"),
        ("ZONE 5", "CONTINUOUS", "Immutable audit trail,\nanomaly detection, compliance\nreporting", LIGHT_BLUE, "\U0001F4CA"),
    ]

    for i, (zone, name, desc, accent, icon) in enumerate(zones):
        y = Inches(2.0 + i * 1.05)
        # Zone number badge
        add_rounded_rect(slide, Inches(0.8), y, Inches(1.4), Inches(0.85),
                         accent, zone, font_size=13, font_color=DARK_NAVY, bold=True)
        # Zone name
        add_text_box(slide, Inches(2.4), y + Inches(0.05), Inches(2.8), Inches(0.4),
                     name, font_size=18, color=WHITE, bold=True)
        # Description
        add_text_box(slide, Inches(5.5), y + Inches(0.0), Inches(5.5), Inches(0.85),
                     desc, font_size=13, color=GRAY, line_spacing=1.3)
        # Connector
        if i < len(zones) - 1:
            add_rounded_rect(slide, Inches(1.3), y + Inches(0.85), Inches(0.3), Inches(0.2),
                             DARK_GRAY)


def slide_05_threat_detection(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "ZONE 1: PRE-MODEL", font_size=14, color=RED, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "Threat Detection Before AI Runs",
                 font_size=32, color=WHITE, bold=True)

    add_divider(slide, Inches(1.6), RED)

    # Big number
    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(3.0), Inches(1.0),
                 "212", font_size=72, color=RED, bold=True)
    add_text_box(slide, Inches(0.8), Inches(2.9), Inches(3.0), Inches(0.4),
                 "Attack patterns scanned", font_size=16, color=GRAY)

    # Categories
    categories = [
        ("Direct Override", "20 patterns", "\"Forget your instructions\""),
        ("SQL Fragments", "31 patterns", "UNION SELECT, DROP TABLE"),
        ("Encoding Bypass", "6 patterns", "Hex, Unicode, URL encoding"),
        ("Indirect Injection", "20 patterns", "Social engineering, urgency"),
        ("Unicode Bypass", "7 patterns", "Zero-width chars, RTL override"),
        ("Prompt Leaking", "10 patterns", "\"Show me your prompt\""),
        ("CoT Manipulation", "8 patterns", "Reasoning-path hijacking"),
        ("Delimiter Injection", "Multi", "ChatML, Llama tag injection"),
    ]

    for i, (cat, count, example) in enumerate(categories):
        col = i % 2
        row = i // 2
        x = Inches(4.5 + col * 4.3)
        y = Inches(1.9 + row * 1.25)
        add_rounded_rect(slide, x, y, Inches(3.8), Inches(1.05), LIGHT_NAVY)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(2.5), Inches(0.35),
                     cat, font_size=14, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(2.8), y + Inches(0.1), Inches(0.8), Inches(0.35),
                     count, font_size=11, color=BLUE, bold=True, alignment=PP_ALIGN.RIGHT)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.5), Inches(3.4), Inches(0.4),
                     example, font_size=11, color=DARK_GRAY)

    # Additional checks
    add_text_box(slide, Inches(0.8), Inches(3.6), Inches(3.2), Inches(0.4),
                 "Also includes:", font_size=14, color=LIGHT_BLUE, bold=True)

    extras = [
        "Schema probing detection (24h sliding window)",
        "Behavioral fingerprinting (30-day user profiles)",
        "Employment status verification (terminated = blocked)",
        "Weighted threat classification (CRITICAL \u2192 NONE)",
    ]
    add_multiline_box(slide, Inches(0.8), Inches(4.1), Inches(3.5), Inches(3.0),
                      extras, font_size=13, color=GRAY, bullet=True, line_spacing=1.6)


def slide_06_ai_guardrails(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "ZONES 2 & 3: MODEL BOUNDARY + POST-MODEL", font_size=14, color=AMBER, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "AI Guardrails: Control, Validate, Rewrite",
                 font_size=32, color=WHITE, bold=True)

    add_divider(slide, Inches(1.6), AMBER)

    # Left: Context Minimization
    add_text_box(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.5),
                 "Context Minimization (Zone 2)", font_size=20, color=AMBER, bold=True)

    ctx_lines = [
        "AI only sees tables & columns the user's role permits",
        "HIDDEN columns completely omitted from schema (zero leakage)",
        "Row-filter rules injected as natural language constraints",
        "Dialect-specific hints (MySQL 8.0 / PostgreSQL 16)",
    ]
    add_multiline_box(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(2.5),
                      ctx_lines, font_size=14, color=GRAY, bullet=True, line_spacing=1.7)

    # Right: 3-Gate Validation
    add_text_box(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.5),
                 "3-Gate SQL Validation (Zone 3)", font_size=20, color=BLUE, bold=True)

    gates = [
        ("Gate 1: Structural", "Table & column authorization,\nsubquery depth, stacked query detection"),
        ("Gate 2: Classification", "Sensitivity levels (L1-L5) vs clearance,\nPII in aggregates, unmasked sensitive data"),
        ("Gate 3: Behavioral", "Write operations, UNION exfiltration,\nsystem table access, privilege escalation"),
    ]
    for i, (gate, desc) in enumerate(gates):
        y = Inches(2.5 + i * 1.15)
        add_rounded_rect(slide, Inches(7.0), y, Inches(5.5), Inches(1.0), LIGHT_NAVY)
        add_text_box(slide, Inches(7.2), y + Inches(0.05), Inches(5.0), Inches(0.35),
                     gate, font_size=14, color=WHITE, bold=True)
        add_text_box(slide, Inches(7.2), y + Inches(0.4), Inches(5.0), Inches(0.5),
                     desc, font_size=12, color=GRAY)

    # Bottom: additional features
    add_rounded_rect(slide, Inches(0.8), Inches(5.8), Inches(3.6), Inches(1.2), LIGHT_NAVY)
    add_text_box(slide, Inches(1.0), Inches(5.9), Inches(3.2), Inches(0.35),
                 "\u2714  Hallucination Detection", font_size=15, color=GREEN, bold=True)
    add_text_box(slide, Inches(1.0), Inches(6.3), Inches(3.2), Inches(0.5),
                 "Catches fabricated table &\ncolumn names in AI output", font_size=12, color=GRAY)

    add_rounded_rect(slide, Inches(4.8), Inches(5.8), Inches(3.6), Inches(1.2), LIGHT_NAVY)
    add_text_box(slide, Inches(5.0), Inches(5.9), Inches(3.2), Inches(0.35),
                 "\u2714  Auto-Masking", font_size=15, color=GREEN, bold=True)
    add_text_box(slide, Inches(5.0), Inches(6.3), Inches(3.2), Inches(0.5),
                 "PII automatically masked:\nfirst initial, year-only, SHA-256", font_size=12, color=GRAY)

    add_rounded_rect(slide, Inches(8.8), Inches(5.8), Inches(3.6), Inches(1.2), LIGHT_NAVY)
    add_text_box(slide, Inches(9.0), Inches(5.9), Inches(3.2), Inches(0.35),
                 "\u2714  Row-Filter Injection", font_size=15, color=GREEN, bold=True)
    add_text_box(slide, Inches(9.0), Inches(6.3), Inches(3.2), Inches(0.5),
                 "Mandatory WHERE clauses\nauto-injected per role policy", font_size=12, color=GRAY)


def slide_07_rbac(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "ACCESS CONTROL", font_size=14, color=BLUE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "Role-Based Access Control (RBAC)",
                 font_size=32, color=WHITE, bold=True)

    add_divider(slide, Inches(1.6))

    # Metrics row
    metrics = [
        ("17", "Roles", BLUE),
        ("5", "Clearance\nLevels", GREEN),
        ("8", "Data\nDomains", AMBER),
        ("3", "Visibility\nModes", RED),
    ]
    for i, (num, label, accent) in enumerate(metrics):
        x = Inches(0.8 + i * 3.1)
        add_text_box(slide, x, Inches(1.9), Inches(1.5), Inches(0.8),
                     num, font_size=48, color=accent, bold=True)
        add_text_box(slide, x + Inches(1.2), Inches(2.0), Inches(1.8), Inches(0.7),
                     label, font_size=14, color=GRAY)

    # Column visibility example
    add_text_box(slide, Inches(0.8), Inches(3.1), Inches(11), Inches(0.5),
                 "Same Query, Different Results:", font_size=18, color=LIGHT_BLUE, bold=True)

    # Role comparison table header
    headers = ["Column", "Nurse (L2)", "Physician (L3)", "Attending (L4)"]
    cols_x = [Inches(0.8), Inches(4.0), Inches(6.8), Inches(9.6)]
    cols_w = [Inches(3.0), Inches(2.6), Inches(2.6), Inches(2.6)]

    y_header = Inches(3.7)
    for j, hdr in enumerate(headers):
        add_rounded_rect(slide, cols_x[j], y_header, cols_w[j], Inches(0.45),
                         BLUE, hdr, font_size=12, font_color=WHITE, bold=True)

    rows = [
        ("patient_id", "VISIBLE", "VISIBLE", "VISIBLE", GREEN, GREEN, GREEN),
        ("first_name", "MASKED  (R***)", "VISIBLE", "VISIBLE", AMBER, GREEN, GREEN),
        ("aadhaar_number", "HIDDEN", "HIDDEN", "VISIBLE", RED, RED, GREEN),
        ("date_of_birth", "HIDDEN", "VISIBLE", "VISIBLE", RED, GREEN, GREEN),
        ("hiv_status", "HIDDEN", "HIDDEN", "HIDDEN", RED, RED, RED),
    ]

    for i, (col, v1, v2, v3, c1, c2, c3) in enumerate(rows):
        y = Inches(4.25 + i * 0.55)
        bg = LIGHT_NAVY if i % 2 == 0 else DARK_NAVY
        add_rounded_rect(slide, cols_x[0], y, cols_w[0], Inches(0.45), bg, col, font_size=12, font_color=WHITE, alignment=PP_ALIGN.LEFT)
        add_rounded_rect(slide, cols_x[1], y, cols_w[1], Inches(0.45), bg, v1, font_size=12, font_color=c1)
        add_rounded_rect(slide, cols_x[2], y, cols_w[2], Inches(0.45), bg, v2, font_size=12, font_color=c2)
        add_rounded_rect(slide, cols_x[3], y, cols_w[3], Inches(0.45), bg, v3, font_size=12, font_color=c3)


def slide_08_multi_db(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "ARCHITECTURE", font_size=14, color=GREEN, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "Multi-Database Architecture",
                 font_size=32, color=WHITE, bold=True)

    add_divider(slide, Inches(1.6), GREEN)

    dbs = [
        ("ApolloHIS", "MySQL 8.0", "12 tables", "Clinical & HIS",
         "patients, encounters, vital_signs,\nlab_results, prescriptions, allergies,\nappointments, clinical_notes...", BLUE),
        ("ApolloHR", "MySQL 8.0", "5 tables", "Human Resources",
         "employees, payroll, leave_records,\ncertifications, credentials", RGBColor(0xEC, 0x48, 0x99)),
        ("apollo_financial", "PostgreSQL 16", "6 tables", "Financial",
         "claims, claim_line_items,\ninsurance_plans, patient_billing,\npayer_contracts, payments", GREEN),
        ("apollo_analytics", "PostgreSQL 16", "4 tables", "Analytics",
         "encounter_summaries,\npopulation_health, quality_metrics,\nresearch_cohorts", LIGHT_BLUE),
    ]

    for i, (name, engine, count, domain, tables, accent) in enumerate(dbs):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 6.3)
        y = Inches(2.0 + row * 2.6)

        add_rounded_rect(slide, x, y, Inches(5.8), Inches(2.3), LIGHT_NAVY)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.15), Inches(3.0), Inches(0.4),
                     name, font_size=20, color=accent, bold=True)
        add_text_box(slide, x + Inches(3.5), y + Inches(0.15), Inches(2.0), Inches(0.4),
                     engine, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.RIGHT)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.55), Inches(3.0), Inches(0.3),
                     f"{count}  |  {domain}", font_size=13, color=GRAY)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.95), Inches(5.2), Inches(1.2),
                     tables, font_size=12, color=DARK_GRAY, line_spacing=1.4)

    # Total badge
    add_rounded_rect(slide, Inches(5.0), Inches(7.0), Inches(3.3), Inches(0.4),
                     BLUE, "27 tables across 4 databases  |  Automatic routing",
                     font_size=12, font_color=WHITE, bold=True)


def slide_09_compliance(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "COMPLIANCE & AUDIT", font_size=14, color=GREEN, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "7 Regulatory Frameworks, One Platform",
                 font_size=32, color=WHITE, bold=True)

    add_divider(slide, Inches(1.6), GREEN)

    frameworks = [
        ("HIPAA Privacy", "45 CFR 164.502-528", "PHI access control, disclosure\naccounting, de-identification"),
        ("HIPAA Security", "45 CFR 164.312", "Access control, audit controls,\nintegrity, authentication"),
        ("42 CFR Part 2", "Substance Use Records", "Hard block on substance abuse\n& psychotherapy records"),
        ("SOX", "Sarbanes-Oxley", "Financial report integrity,\nimmutable audit trail"),
        ("GDPR", "EU Data Protection", "Data protection by design,\nprocessing records"),
        ("EU AI Act", "High-Risk AI Systems", "Risk management, transparency,\nhuman oversight"),
        ("ISO 42001", "AI Management Systems", "Impact assessment, monitoring,\ncontinual improvement"),
    ]

    for i, (name, subtitle, desc) in enumerate(frameworks):
        col = i % 4
        row = i // 4
        x = Inches(0.8 + col * 3.1)
        y = Inches(2.0 + row * 2.5)

        add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.2), LIGHT_NAVY)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.4), Inches(0.4),
                     name, font_size=16, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.55), Inches(2.4), Inches(0.3),
                     subtitle, font_size=10, color=BLUE)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.95), Inches(2.4), Inches(1.0),
                     desc, font_size=12, color=GRAY, line_spacing=1.4)

    # Audit trail highlight
    add_rounded_rect(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.9), DARK_NAVY)
    add_text_box(slide, Inches(1.0), Inches(6.3), Inches(5.0), Inches(0.3),
                 "SHA-256 Hash-Chain Audit Trail", font_size=16, color=GREEN, bold=True)
    add_text_box(slide, Inches(1.0), Inches(6.65), Inches(5.0), Inches(0.3),
                 "Tamper-detectable  |  Append-only  |  Every query logged", font_size=12, color=GRAY)
    add_text_box(slide, Inches(8.5), Inches(6.3), Inches(3.5), Inches(0.3),
                 "< 3 seconds", font_size=28, color=BLUE, bold=True, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(8.5), Inches(6.7), Inches(3.5), Inches(0.3),
                 "per compliance report", font_size=12, color=GRAY, alignment=PP_ALIGN.RIGHT)


def slide_10_btg(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "EMERGENCY ACCESS", font_size=14, color=RED, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "Break-the-Glass: Controlled Emergency Override",
                 font_size=32, color=WHITE, bold=True)

    add_divider(slide, Inches(1.6), RED)

    # How it works
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.5),
                 "How It Works", font_size=20, color=LIGHT_BLUE, bold=True)

    steps = [
        ("1", "Emergency Declared", "Authorized role initiates BTG\nwith mandatory reason"),
        ("2", "4-Hour Token Issued", "Clearance temporarily elevated\nwith strict time limit"),
        ("3", "Full Audit Trail", "Every action logged, compliance\nofficer notified immediately"),
        ("4", "Auto-Expiration", "Access reverts after 4 hours\nwith mandatory justification"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        y = Inches(2.7 + i * 1.1)
        add_rounded_rect(slide, Inches(0.8), y, Inches(0.6), Inches(0.6),
                         BLUE, num, font_size=18, font_color=WHITE, bold=True)
        add_text_box(slide, Inches(1.6), y + Inches(0.0), Inches(2.5), Inches(0.35),
                     title, font_size=16, color=WHITE, bold=True)
        add_text_box(slide, Inches(1.6), y + Inches(0.35), Inches(3.5), Inches(0.6),
                     desc, font_size=12, color=GRAY, line_spacing=1.3)

    # Hard limits
    add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5.0), Inches(0.5),
                 "Hard Limits (Never Overridable)", font_size=20, color=RED, bold=True)

    add_rounded_rect(slide, Inches(7.0), Inches(2.7), Inches(5.2), Inches(4.0), LIGHT_NAVY)

    hard_limits = [
        "Psychotherapy notes (42 CFR Part 2)",
        "Substance abuse records (42 CFR Part 2)",
        "HIV status (Federal & State law)",
        "Genetic testing data (GINA Act)",
    ]

    add_text_box(slide, Inches(7.3), Inches(2.9), Inches(4.6), Inches(0.4),
                 "Even in emergency mode, these are ALWAYS blocked:",
                 font_size=13, color=AMBER)

    for i, limit in enumerate(hard_limits):
        y = Inches(3.5 + i * 0.7)
        add_text_box(slide, Inches(7.3), y, Inches(4.6), Inches(0.4),
                     f"\u2718  {limit}", font_size=15, color=RED, bold=True)

    add_text_box(slide, Inches(7.3), Inches(6.0), Inches(4.6), Inches(0.5),
                 "No exceptions. No overrides.\nPatient privacy is absolute.",
                 font_size=13, color=GRAY, line_spacing=1.4)


def slide_11_metrics(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                 "BY THE NUMBERS", font_size=14, color=BLUE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                 "QueryVault at a Glance",
                 font_size=32, color=WHITE, bold=True)

    add_divider(slide, Inches(1.6))

    metrics = [
        ("212", "Attack Patterns\nScanned", RED),
        ("5", "Security\nZones", BLUE),
        ("17", "RBAC\nRoles", GREEN),
        ("5", "Clearance\nLevels", AMBER),
        ("8", "Data\nDomains", LIGHT_BLUE),
        ("27", "Database\nTables", GREEN),
        ("7", "Compliance\nFrameworks", BLUE),
        ("< 3s", "Report\nGeneration", AMBER),
        ("100%", "Audit\nCoverage", GREEN),
        ("4h", "BTG Time\nLimit", RED),
    ]

    for i, (num, label, accent) in enumerate(metrics):
        col = i % 5
        row = i // 5
        x = Inches(0.5 + col * 2.5)
        y = Inches(2.2 + row * 2.6)
        add_metric_card(slide, x, y, num, label, accent)


def slide_12_thankyou(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_NAVY)

    add_divider(slide, Inches(3.0), BLUE, Inches(4.5), Inches(4.3))

    add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.8),
                 "Thank You", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(0.6),
                 "Ready to secure your AI-powered database queries?",
                 font_size=20, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

    # CTA cards
    ctas = [
        ("\U0001F4BB", "Live Demo", "See QueryVault in action\nwith real hospital scenarios"),
        ("\U0001F4DE", "Contact Us", "Schedule a consultation\nwith our security team"),
        ("\U0001F4C4", "Technical Brief", "Deep-dive into our\n5-zone architecture"),
    ]

    for i, (icon, title, desc) in enumerate(ctas):
        x = Inches(1.5 + i * 3.8)
        add_rounded_rect(slide, x, Inches(4.3), Inches(3.2), Inches(2.0), LIGHT_NAVY)
        add_text_box(slide, x + Inches(0.3), Inches(4.4), Inches(2.6), Inches(0.5),
                     icon, font_size=28, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.3), Inches(5.0), Inches(2.6), Inches(0.4),
                     title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.3), Inches(5.4), Inches(2.6), Inches(0.7),
                     desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    add_text_box(slide, Inches(1.5), Inches(6.8), Inches(10), Inches(0.4),
                 "QueryVault  |  Apollo Hospitals  |  AI Query Security Platform",
                 font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# ── Main ────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_pipeline(prs)
    slide_05_threat_detection(prs)
    slide_06_ai_guardrails(prs)
    slide_07_rbac(prs)
    slide_08_multi_db(prs)
    slide_09_compliance(prs)
    slide_10_btg(prs)
    slide_11_metrics(prs)
    slide_12_thankyou(prs)

    output = "QueryVault_Sales_Deck.pptx"
    prs.save(output)
    print(f"\n  Sales deck generated: {output}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Open in PowerPoint, Keynote, or Google Slides\n")


if __name__ == "__main__":
    main()
